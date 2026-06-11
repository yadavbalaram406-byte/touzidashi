import base64
import subprocess
import tempfile
from pathlib import Path

import httpx
import pdfplumber


def _pdf_page_to_base64(file_path: str, page_index: int) -> str:
    import fitz
    doc = fitz.open(file_path)
    pix = doc[page_index].get_pixmap(matrix=fitz.Matrix(1.5, 1.5))
    b64 = base64.b64encode(pix.tobytes("png")).decode()
    doc.close()
    return b64


async def _vision_extract_pdf(file_path: str) -> str:
    """Extract text from scanned PDF via Vision API (page by page)."""
    import fitz
    from app.config import get_settings

    settings = get_settings()
    doc = fitz.open(file_path)
    total_pages = min(len(doc), 20)
    doc.close()

    results = []
    for i in range(total_pages):
        img_b64 = _pdf_page_to_base64(file_path, i)
        payload = {
            "model": settings.anthropic_model,
            "max_tokens": 2000,
            "messages": [{
                "role": "user",
                "content": [
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_b64}"}},
                    {"type": "text", "text": (
                        "请将这张图片中所有可见的文字内容完整转录出来，保持原始段落和列表结构。"
                        "如果是图表或表格，用文字简要说明内容。只输出转录文字，不要添加解释。"
                    )},
                ],
            }],
        }
        base_url = settings.anthropic_base_url.rstrip("/")
        async with httpx.AsyncClient(timeout=60) as client:
            r = await client.post(
                f"{base_url}/v1/chat/completions",
                headers={"Authorization": f"Bearer {settings.anthropic_api_key}", "Content-Type": "application/json"},
                json=payload,
            )
            r.raise_for_status()
            text = r.json()["choices"][0]["message"]["content"]
            if text.strip():
                results.append(f"[第{i+1}页]\n{text.strip()}")

    return "\n\n".join(results)


async def extract_text_from_pdf(file_path: str) -> str:
    # Fast path: try text extraction first
    text_parts = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
    result = "\n\n".join(text_parts)

    # Scanned PDF: fall back to Vision API
    if not result.strip():
        result = await _vision_extract_pdf(file_path)

    return result


def extract_text_from_pptx(file_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(file_path)
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
        if parts:
            slides_text.append(f"[第{i}页]\n" + "\n".join(parts))
    return "\n\n".join(slides_text)


def _convert_ppt_to_pdf(file_path: str) -> str | None:
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", tmpdir, file_path],
                capture_output=True,
                timeout=60,
            )
            if result.returncode == 0:
                stem = Path(file_path).stem
                out = Path(tmpdir) / f"{stem}.pdf"
                if out.exists():
                    return str(out)
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass
    return None


async def extract_text_from_ppt(file_path: str) -> str:
    pdf_path = _convert_ppt_to_pdf(file_path)
    if pdf_path:
        return await extract_text_from_pdf(pdf_path)
    try:
        return extract_text_from_pptx(file_path)
    except Exception:
        return "[警告：.ppt格式需安装LibreOffice才能完整解析，文本提取可能不完整]"


def truncate_text(text: str, max_chars: int = 100_000) -> str:
    if len(text) <= max_chars:
        return text
    keep_start = 80_000
    keep_end = 20_000
    return text[:keep_start] + f"\n\n[...内容过长，已省略中间部分...]\n\n" + text[-keep_end:]


async def parse_document(file_path: str, file_type: str) -> tuple[str, int]:
    """Returns (extracted_text, text_length)."""
    ext = file_type.lower().lstrip(".")
    if ext == "pdf":
        raw = await extract_text_from_pdf(file_path)
    elif ext == "pptx":
        raw = extract_text_from_pptx(file_path)
    elif ext == "ppt":
        raw = await extract_text_from_ppt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    text = truncate_text(raw.strip())
    return text, len(text)
